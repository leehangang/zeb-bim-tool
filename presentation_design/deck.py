# -*- coding: utf-8 -*-
"""
ZEB-ROI 졸업설계 발표 덱 — MBB 스타일, Noto Sans KR, 인포그래픽 중심.
Pillow 정밀 렌더(PNG) → python-pptx 조립(.pptx).
"""
import os, math
from PIL import Image, ImageDraw, ImageFont

NOTO = "C:/Windows/Fonts/NotoSansKR-VF.ttf"
BASE = "C:/Users/이혁주/Desktop/zeb-chatbot/presentation_design/"
SLIDES = BASE + "slides/"
os.makedirs(SLIDES, exist_ok=True)

W, H = 1280, 720
SCALE = 3
WHITE=(255,255,255); INK=(26,32,28); GREEN=(27,94,32); GREEN_D=(17,53,27)
GREEN_DK=(12,40,21); GOLD=(193,138,45); GRAY=(99,110,104); GRAY_L=(155,163,158)
LINE=(226,228,224); TINT=(238,243,239); TINT2=(244,239,228)
MX=88; TITLE_SZ=31

_fc={}
def F(weight,size):
    key=(weight,round(size*SCALE))
    if key in _fc: return _fc[key]
    f=ImageFont.truetype(NOTO,round(size*SCALE))
    for arg in (weight,weight.encode()):
        try: f.set_variation_by_name(arg); break
        except Exception: continue
    _fc[key]=f; return f

class Slide:
    def __init__(s,bg=WHITE):
        s.img=Image.new("RGB",(W*SCALE,H*SCALE),bg); s.d=ImageDraw.Draw(s.img)
    def T(s,x,y,t,w,sz,c,anchor="la",track=0.0):
        f=F(w,sz)
        if track==0: s.d.text((x*SCALE,y*SCALE),t,font=f,fill=c,anchor=anchor)
        else:
            cx=x*SCALE
            for ch in t:
                s.d.text((cx,y*SCALE),ch,font=f,fill=c,anchor=anchor); cx+=f.getlength(ch)+track*SCALE
    def line(s,x0,y0,x1,y1,c,w=1.0): s.d.line([x0*SCALE,y0*SCALE,x1*SCALE,y1*SCALE],fill=c,width=max(1,round(w*SCALE)))
    def rect(s,x0,y0,x1,y1,fill=None,outline=None,w=1.0,r=0):
        bb=[x0*SCALE,y0*SCALE,x1*SCALE,y1*SCALE]
        if r>0: s.d.rounded_rectangle(bb,radius=r*SCALE,fill=fill,outline=outline,width=max(1,round(w*SCALE)))
        else: s.d.rectangle(bb,fill=fill,outline=outline,width=max(1,round(w*SCALE)))
    def circle(s,cx,cy,r,fill=None,outline=None,w=1.0):
        s.d.ellipse([(cx-r)*SCALE,(cy-r)*SCALE,(cx+r)*SCALE,(cy+r)*SCALE],fill=fill,outline=outline,width=max(1,round(w*SCALE)))
    def poly(s,pts,fill): s.d.polygon([(x*SCALE,y*SCALE) for x,y in pts],fill=fill)
    def measure(s,t,w,sz): return F(w,sz).getlength(t)/SCALE
    def lh(s,w,sz): a,dd=F(w,sz).getmetrics(); return (a+dd)/SCALE
    def wrap(s,t,w,sz,maxw):
        out=[]; cur=""
        for wd in t.split(" "):
            tt=wd if not cur else cur+" "+wd
            if s.measure(tt,w,sz)<=maxw: cur=tt
            else:
                if cur: out.append(cur)
                cur=wd
        if cur: out.append(cur)
        return out
    def para(s,x,y,t,w,sz,c,maxw,leading=1.5):
        l=s.lh(w,sz)*leading
        for i,ln in enumerate(s.wrap(t,w,sz,maxw)): s.T(x,y+i*l,ln,w,sz,c)
        return y+len(s.wrap(t,w,sz,maxw))*l
    def save(s,n): s.img.save(SLIDES+n,"PNG")

# furniture ----------------------------------------------------------------
def title_block(s,eyebrow,title,dark=False):
    s.T(MX,62,eyebrow,"Bold",14,GOLD,track=2.0)
    y=90; tc=WHITE if dark else INK
    for ln in s.wrap(title,"Bold",TITLE_SZ,W-2*MX-40):
        s.T(MX,y,ln,"Bold",TITLE_SZ,tc); y+=s.lh("Bold",TITLE_SZ)*1.12
    return y
def footer(s,page,source="출처: 그린리모델링 가이드라인 · 조달청 단가DB · ZEB 인증기준 고시",dark=False):
    c=(140,156,146) if dark else GRAY_L
    s.line(MX,678,W-MX,678,(58,78,64) if dark else LINE,0.8)
    s.T(MX,690,f"{page:02d}","Medium",11,c)
    s.T(MX+34,690,"ZEB-ROI · 그린리모델링 의사결정 플랫폼","Regular",11,c)
    s.T(W-MX,690,source,"Regular",10.5,c,anchor="ra")
def token(s,x,y,num,color=GREEN,r=17):
    s.circle(x,y,r,fill=color); s.T(x,y-1,num,"Bold",15,WHITE,anchor="mm")
def arrow(s,x0,x1,y,color=GRAY_L,w=1.6):
    s.line(x0,y,x1-7,y,color,w); s.poly([(x1,y),(x1-9,y-5),(x1-9,y+5)],color)
def chip(s,x,y,txt,fg,bg,sz=12.5,h=28,padx=13):
    w=s.measure(txt,"Medium",sz)+2*padx; s.rect(x,y,x+w,y+h,fill=bg,r=h/2)
    s.T(x+w/2,y+h/2-1,txt,"Medium",sz,fg,anchor="mm"); return w

# --- 벡터 아이콘/일러스트 (직접 그림, 저작권 무관) -----------------------
def i_sun(s,cx,cy,r,color):
    s.circle(cx,cy,r,outline=color,w=2.2)
    for k in range(8):
        a=math.radians(k*45)
        s.line(cx+(r+3)*math.cos(a),cy+(r+3)*math.sin(a),cx+(r+9)*math.cos(a),cy+(r+9)*math.sin(a),color,1.8)
def i_bim(s,cx,cy,sz,color):
    o=sz*0.3
    s.rect(cx-sz/2+o,cy-sz/2-o,cx+sz/2+o,cy+sz/2-o,outline=color,w=1.4)
    s.rect(cx-sz/2,cy-sz/2,cx+sz/2,cy+sz/2,outline=color,w=2.2)
    s.line(cx,cy-sz/2,cx,cy+sz/2,color,1.2); s.line(cx-sz/2,cy,cx+sz/2,cy,color,1.2)
def skyline(s,x0,base,col):
    # 공공건축물 스카이라인 실루엣 (subtle background motif)
    specs=[(118,72,"flat"),(86,104,"pitch"),(150,58,"win"),(74,118,"tower"),
           (132,84,"solar"),(96,96,"pitch"),(140,64,"win")]
    x=x0
    for w,h,t in specs:
        top=base-h
        s.rect(x,top,x+w,base,fill=col)
        if t=="pitch":
            s.poly([(x-6,top),(x+w/2,top-26),(x+w+6,top)],col)
        if t=="solar":
            s.rect(x+18,top-14,x+w-18,top-3,fill=(150,176,156))
            for k in range(4): s.line(x+18+(k+1)*(w-36)/5,top-14,x+18+(k+1)*(w-36)/5,top-3,col,1.0)
        if t=="win":
            for r in range(2):
                for c in range(int(w//34)):
                    s.rect(x+14+c*34,top+16+r*26,x+14+c*34+18,top+16+r*26+14,fill=(255,255,255))
        if t=="tower":
            for r in range(4): s.rect(x+18,top+14+r*24,x+w-18,top+14+r*24+12,fill=(255,255,255))
        x+=w+14

# ==========================================================================
def s01():
    s=Slide(GREEN_DK)
    for i in range(34):
        t=i/33; x=720+t*(W-MX-720); y=560-(t**1.7)*300; a=int(40+t*120)
        s.circle(x,y,2.0,fill=(a,int(a*0.8)+30,70))
    s.T(MX,150,"ZERO ENERGY BUILDING · GREEN REMODELING","Medium",14,(150,200,160),track=2.2)
    s.T(MX,196,"BIM 한 번으로 ZEB 등급을 평가하고","Bold",46,WHITE)
    s.T(MX,254,"그린리모델링 전 과정을 설계하다","Bold",46,WHITE)
    s.T(MX,340,"제로에너지건축물(ZEB) 의사결정 통합 플랫폼","Regular",19,(190,205,193))
    s.line(MX,604,MX+360,604,(70,95,76),1.0)
    s.T(MX,624,"졸업설계 작품 · 2026","Medium",14,(175,190,178))
    s.T(W-MX,624,"사례 · 공공기관 소유 어린이집","Medium",13,GOLD,anchor="ra")
    s.save("slide-01.png")

def s02():
    s=Slide(WHITE)
    title_block(s,"AGENDA","발표 순서")
    items=[("01","배경","ZEB 의무화와 공공건축 BIM 의무화"),
           ("02","문제","ZEB 의사결정과 현행 제도의 한계"),
           ("03","제안","BIM 기반 통합 의사결정 플랫폼"),
           ("04","실증","공공기관 소유 어린이집 검증"),
           ("05","전망","기대효과와 앞으로의 AI 발전방향")]
    y=212; rh=88
    for num,t,d in items:
        token(s,MX+18,y+18,num,GREEN,19)
        s.T(MX+58,y,t,"Bold",22,INK)
        s.T(MX+58,y+34,d,"Regular",15,GRAY)
        if num!="05": s.line(MX+58,y+70,W-MX,y+70,LINE,0.8)
        y+=rh
    footer(s,2)
    s.save("slide-02.png")

def s03():
    s=Slide(WHITE)
    title_block(s,"01 · 배경","건물부문 탄소중립을 위해 ZEB는 공공에서 민간으로 의무화가 확대된다")
    # roadmap timeline
    y=320; x0=MX+30; x1=W-MX-30
    s.line(x0,y,x1,y,LINE,2.0)
    nodes=[(0.04,"2020","공공 1,000㎡ 이상\nZEB 5등급 의무화",GREEN),
           (0.40,"단계 확대","공공 연면적 기준\n점진적 강화",GREEN),
           (0.74,"민간 확대","민간 건축물로\nZEB 의무 확대",GOLD),
           (0.98,"기존 건축물","그린리모델링이\n핵심 과제로",GOLD)]
    for t,lab,desc,col in nodes:
        x=x0+t*(x1-x0)
        s.circle(x,y,8,fill=col)
        s.T(x,y-58,lab,"Bold",17,INK,anchor="ma")
        for i,dl in enumerate(desc.split("\n")):
            s.T(x,y+28+i*22,dl,"Regular",13,GRAY,anchor="ma")
    s.rect(MX,470,W-MX,556,fill=TINT,r=10)
    s.T(MX+26,492,"핵심 시사점","Bold",14,GREEN)
    s.T(MX+26,518,"신축뿐 아니라 이미 지어진 건축물의 에너지 성능 개선(그린리모델링)이 ZEB 달성의 관건이다.","Regular",15,INK)
    s.T(MX,600,"* ZEB 의무화 연도·등급은 최신 국토교통부 로드맵 기준으로 확인 필요","Regular",11.5,GRAY_L)
    footer(s,3,"출처: 국토교통부 제로에너지건축물 인증 의무화 로드맵")
    s.save("slide-03.png")

def s04():
    s=Slide(WHITE)
    title_block(s,"01 · 배경","공공 발주 BIM 의무화로, 건물 데이터는 이미 디지털로 쌓이고 있다")
    # left concept
    s.para(MX,232,"조달청·공공 발주 사업에서 BIM(건물정보모델) 적용이 단계적으로 의무화되며, 설계 단계의 건물 정보가 디지털 객체로 축적되고 있다.","Regular",16,INK,520)
    s.rect(MX,360,MX+520,470,fill=TINT,r=10)
    s.T(MX+26,384,"그런데","Bold",14,GOLD)
    s.para(MX+26,410,"쌓인 BIM 데이터는 대부분 설계·시공 관리에만 쓰이고, 에너지·경제성 의사결정으로는 거의 활용되지 못한다.","Regular",15,INK,470)
    # right: BIM object data
    rx=MX+560
    s.T(rx,232,"BIM 객체에 이미 담긴 정량 데이터","Bold",16,INK)
    data=[("벽 · 지붕 · 바닥","단열 적용 면적 / 미적용 면적"),
          ("창 · 문","면적 · 사양 · 열관류율"),
          ("설비","냉·난방·급탕 용량"),
          ("신재생","태양광 용량 → 자립률")]
    yy=274
    for a,b in data:
        s.circle(rx+8,yy+11,4,fill=GREEN)
        s.T(rx+28,yy,a,"Medium",15,INK)
        s.T(rx+28,yy+22,b,"Regular",13,GRAY)
        yy+=64
    arrow(s,rx+150,rx+260,yy+4,GOLD,1.6)
    s.T(rx+270,yy-7,"진단 입력으로 활용","Bold",14,GOLD)
    footer(s,4,"출처: 조달청 BIM 적용 의무화 지침")
    s.save("slide-04.png")

def s05():
    s=Slide(WHITE)
    title_block(s,"01 · 배경","두 의무화가 만나는 지점에 '의사결정 도구'가 비어 있다")
    s.T(MX,232,"ZEB는 결과(등급)를 요구하고, BIM은 데이터(입력)를 제공한다 — 그 사이를 잇는 통합 판단 도구가 없다.","Regular",16,GRAY)
    cy=410
    # ZEB box (left)
    s.rect(MX,318,MX+320,500,fill=TINT,r=12)
    i_sun(s,MX+160,360,16,GREEN)
    s.T(MX+160,392,"ZEB 의무화","Bold",20,GREEN,anchor="ma")
    s.T(MX+160,430,"건물은 정해진 ZEB 등급을 달성해야 한다","Regular",13.5,GRAY,anchor="ma")
    # BIM box (right)
    bx=W-MX-320
    s.rect(bx,318,bx+320,500,fill=TINT,r=12)
    i_bim(s,bx+160,360,30,GREEN)
    s.T(bx+160,392,"BIM 의무화","Bold",20,GREEN,anchor="ma")
    s.T(bx+160,430,"건물 데이터는 디지털로 쌓이고 있다","Regular",13.5,GRAY,anchor="ma")
    # center gap
    arrow(s,MX+330,W/2-150,cy,GRAY_L,1.8)
    arrow(s,W-MX-330,W/2+150,cy,GRAY_L,1.8)
    s.circle(W/2,cy,82,outline=GOLD,w=2.0)
    s.T(W/2,cy-16,"의사결정","Bold",18,INK,anchor="mm")
    s.T(W/2,cy+12,"도구 부재","Bold",18,GOLD,anchor="mm")
    s.T(W/2,572,"\"ZEB는 의무인데, BIM 데이터로 등급·비용·회수를 통합 판단할 도구가 없다\"","Bold",17,INK,anchor="ma")
    footer(s,5)
    s.save("slide-05.png")

def s06():
    s=Slide(WHITE)
    title_block(s,"02 · 문제","ZEB 의사결정은 4개 분야로 쪼개져 수일~수주가 걸린다")
    cols=[("01","건축사","11개 기술요소\n진단표 수기 작성"),
          ("02","설비기사","에너지 성능 ·\n자립률 계산"),
          ("03","세무사","보조금 · 취득세 ·\n용적률 검토"),
          ("04","시공사","항목별 견적\n(수일~수주)")]
    cw=(W-2*MX-3*22)/4
    for i,(n,t,d) in enumerate(cols):
        x=MX+i*(cw+22)
        s.rect(x,234,x+cw,420,fill=TINT,r=10)
        token(s,x+34,268,n,GREEN,17)
        s.T(x+24,300,t,"Bold",19,INK)
        for j,dl in enumerate(d.split("\n")): s.T(x+24,334+j*24,dl,"Regular",13.5,GRAY)
    arrow(s,MX+40,W-MX-40,470,GRAY_L,1.4)
    s.T(W/2,486,"분야별로 따로 계산 → 통합이 없어 사업 판단이 지연되고, 경제성 근거가 부족하다","Bold",16,INK,anchor="ma")
    s.rect(MX,540,W-MX,604,fill=TINT2,r=10)
    s.T(MX+26,560,"인식의 벽","Bold",14,GOLD)
    s.T(MX+26,584,"\"그린리모델링은 회수기간이 30~50년\"이라는 통념이 ZEB 사업 판단을 가로막는다.","Regular",14.5,INK)
    footer(s,6)
    s.save("slide-06.png")

def s07():
    s=Slide(WHITE)
    title_block(s,"02 · 문제","신청은 수기, 정보는 흩어져 있어 사업 판단 자체가 어렵다")
    rows=[("01","서류·수기 중심의 신청 절차","그린리모델링·ZEB 인증 신청이 복잡한 서식 중심이라 진입장벽이 높다"),
          ("02","흩어진 인센티브 정보","보조금·취득세·용적률 완화 근거가 여러 법·고시에 분산돼 파악이 어렵다"),
          ("03","통합 산정 기준의 부재","비용·보조금·ZEB 등급을 한 번에 산정할 표준 도구가 제도적으로 없다")]
    y=240
    for n,t,d in rows:
        token(s,MX+18,y+16,n,GREEN,18)
        s.T(MX+56,y,t,"Bold",19,INK)
        s.para(MX+56,y+30,d,"Regular",15,GRAY,W-MX-56-MX)
        if n!="03": s.line(MX+56,y+78,W-MX,y+78,LINE,0.8)
        y+=104
    s.rect(MX,548,W-MX,616,fill=TINT,r=10)
    s.T(MX+26,568,"결과","Bold",14,GOLD)
    s.T(MX+26,592,"ZEB 사업의 진입장벽이 높아지고, 등급·비용·회수를 통합 판단할 객관적 근거가 부족하다.","Regular",14.5,INK)
    footer(s,7,"출처: 04 녹색건축법 · 05 지방세특례제한법 · 01 그린리모델링 가이드라인")
    s.save("slide-07.png")

def s08():
    s=Slide(WHITE)
    title_block(s,"02 · 문제 → 개선","통합·자동화·근거기반으로 의사결정을 한 자리에 모은다")
    pill=[("통합","BIM 단일 입력으로 진단부터 경제성까지 하나의 흐름으로"),
          ("자동화","자연어·BIM에서 입력 조건을 자동 추출해 즉시 산출"),
          ("근거기반","단가DB·법령 원문을 인용해 결과의 신뢰성을 확보")]
    cw=(W-2*MX-2*26)/3
    for i,(t,d) in enumerate(pill):
        x=MX+i*(cw+26)
        s.T(x,250,f"0{i+1}","Bold",30,GOLD)
        s.line(x,300,x+cw-20,300,GREEN,2.2)
        s.T(x,320,t,"Bold",22,INK)
        s.para(x,360,d,"Regular",15,GRAY,cw-20)
    s.rect(MX,520,W-MX,588,fill=TINT,r=10)
    s.T(MX+26,540,"개선 방향 한 줄","Bold",14,GREEN)
    s.T(MX+26,564,"분야별로 흩어진 ZEB 의사결정을, BIM 한 번으로 통합·자동화·근거기반화한다.","Regular",15,INK)
    footer(s,8)
    s.save("slide-08.png")

def s09():
    s=Slide(GREEN_DK)
    s.T(MX,150,"03 · 제안","Bold",14,GOLD,track=2.0)
    s.T(MX,196,"BIM 한 번으로","Light",30,(190,205,193))
    s.T(MX,242,"ZEB 등급을 평가하고,","Bold",40,WHITE)
    s.T(MX,300,"그 등급에 도달하는 그린리모델링","Bold",40,WHITE)
    s.T(MX,358,"전 과정을 자동으로 산출한다","Bold",40,WHITE)
    w1=chip(s,MX,452,"그린리모델링 · 사업 · 비용 · 신청",(180,220,185),(30,66,38),13,34,16)
    chip(s,MX+w1+20,452,"ZEB · 등급 · 에너지 성능",(225,200,140),(70,55,24),13,34,16)
    s.T(MX,540,"하나의 입력으로 ZEB(등급)와 그린리모델링(사업)을 분리해 명확히 다룬다.","Regular",16,(180,195,184))
    s.save("slide-09.png")

def s10():
    s=Slide(WHITE)
    title_block(s,"03 · 제안","BIM 단일 입력이 진단부터 신청서까지 하나의 흐름으로 이어진다")
    zy=250; zh=300
    # three zones
    def zone(x,w,head,items,fill):
        s.rect(x,zy,x+w,zy+zh,fill=fill,r=12)
        s.T(x+24,zy+24,head,"Bold",16,GREEN)
        yy=zy+62
        for it in items:
            s.circle(x+30,yy+9,3.5,fill=GREEN); s.T(x+46,yy,it,"Regular",14.5,INK); yy+=38
    z1=MX; w1=300
    z2=MX+340; w2=300
    z3=MX+680; w3=W-MX-(MX+680)
    zone(z1,w1,"입력",["Revit BIM (JSON)","자연어 질문 / 조건"],TINT)
    zone(z2,w2,"처리 엔진",["BIM 진단 + ZEB 평가","ROI 경제성 산정","정책 RAG 검색","사업 신청서 생성"],TINT)
    zone(z3,w3,"출력",["ZEB 인증 등급","보강 우선순위 · Max Cost","보조금 · 회수기간","사업 신청서 초안"],TINT)
    arrow(s,z1+w1,z2,zy+zh/2,GRAY_L,1.8)
    arrow(s,z2+w2,z3,zy+zh/2,GRAY_L,1.8)
    s.T(W/2,600,"분야별 수일~수주의 계산을 BIM 한 번으로 통합한다","Bold",16,INK,anchor="ma")
    footer(s,10)
    s.save("slide-10.png")

def s11():
    s=Slide(WHITE)
    title_block(s,"03 · 제안","하나의 플랫폼이 6개 모드로 진단·정책·경제성·신청·근거를 다룬다")
    modes=[("01","BIM 진단 + ZEB 등급 평가","BIM 업로드 → 11개 기술요소 진단 → ZEB 인증 등급 자동 산정","핵심"),
           ("02","정책 Q&A","11개 법령·고시·공고 원문에서 근거 조항을 인용해 답변",None),
           ("03","ROI 시뮬레이션","자연어 조건 → 공사비·보조금·회수기간 즉시 계산",None),
           ("04","사업 신청 인테이크","대화로 정보를 모아 사업 신청서 초안을 자동 생성",None)]
    cw=(W-2*MX-30)/2; ch=150
    for i,(n,t,d,tag) in enumerate(modes):
        x=MX+(i%2)*(cw+30); y=232+(i//2)*(ch+24)
        s.rect(x,y,x+cw,y+ch,fill=TINT,r=12)
        token(s,x+34,y+38,n,GREEN,18)
        s.T(x+66,y+22,t,"Bold",18,INK)
        if tag: chip(s,x+cw-86,y+22,tag,WHITE,GOLD,12,26,12)
        s.para(x+66,y+58,d,"Regular",14,GRAY,cw-90)
    footer(s,11)
    s.save("slide-11.png")

def s12():
    s=Slide(WHITE)
    title_block(s,"03 · 핵심 방법론","ZEB 인증은 자립률 또는 1차에너지소요량, 그리고 BEMS로 결정된다")
    y=250
    b1=(MX,"제1호","에너지 자립률"); b2=(MX+300,"제2호","1차에너지소요량")
    for x,a,bb in (b1,b2):
        s.rect(x,y,x+250,y+86,fill=TINT,r=10)
        s.T(x+24,y+20,a,"Bold",15,GREEN); s.T(x+24,y+46,bb,"Medium",16,INK)
    s.T(MX+275,y+34,"OR","Bold",16,GOLD,anchor="mm")
    s.T(MX+600,y+30,"AND","Bold",16,GOLD,anchor="lm")
    s.rect(MX+660,y,MX+910,y+86,fill=TINT,r=10)
    s.T(MX+684,y+20,"제3호","Bold",15,GREEN); s.T(MX+684,y+46,"BEMS 설치","Medium",16,INK)
    arrow(s,MX+930,W-MX-150,y+43,GRAY_L,1.8)
    s.T(W-MX,y+30,"인증 등급","Bold",18,INK,anchor="ra")
    # ladder
    ly=400
    s.T(MX,ly,"ZEB 등급 — 에너지 자립률 기준(제1호)","Bold",15,INK)
    grades=[("+등급","120%"),("1등급","100%"),("2등급","80%"),("3등급","60%"),("4등급","40%"),("5등급","20%")]
    cw=(W-2*MX-5*16)/6
    for i,(g,v) in enumerate(grades):
        x=MX+i*(cw+16); col=GREEN if i<3 else GOLD
        s.rect(x,ly+34,x+cw,ly+120,fill=TINT,r=8)
        s.T(x+cw/2,ly+58,g,"Bold",17,col,anchor="ma")
        s.T(x+cw/2,ly+88,v,"Medium",15,GRAY,anchor="ma")
    s.rect(MX,ly+150,W-MX,ly+150+50,outline=None)
    s.T(MX,ly+158,"전력 1차에너지 환산계수 × 2.75 적용  ·  인증등급은 제1·2호 중 상위 등급으로 산정","Regular",13.5,GRAY)
    footer(s,12,"출처: 03 ZEB 인증기준 고시 · 건축물 에너지효율등급 기준")
    s.save("slide-12.png")

def s13():
    s=Slide(WHITE)
    title_block(s,"04 · 실증 사례","검증 사례 — 공공기관 소유 어린이집")
    facts=[("1,251㎡","연면적"),("2014년","사용승인"),("어린이집","건물 용도"),("중부2","기후대")]
    cw=(W-2*MX-3*22)/4
    for i,(v,l) in enumerate(facts):
        x=MX+i*(cw+22)
        s.rect(x,206,x+cw,322,fill=TINT,r=10)
        s.T(x+24,238,v,"Bold",31,GREEN); s.T(x+24,288,l,"Regular",14,GRAY)
    s.line(MX,356,W-MX,356,LINE,0.8)
    s.T(MX,386,"현재 상태","Bold",16,INK)
    s.para(MX,424,"태양광·바닥 단열 등 일부만 보강된 상태로, 외피·설비 보강 여지가 크다. 어린이집 용도라 일조·채광 등 추가 기준도 적용된다.","Regular",15.5,GRAY,520)
    rx=MX+600
    s.T(rx,386,"왜 이 사례인가","Bold",16,INK)
    for i,t in enumerate(["보강 잠재력이 큰 노후 공공건축물","실측 BIM 기반 검증 케이스","ZEB·그린리모델링 의사결정의 전형"]):
        s.circle(rx+8,422+i*42+9,4,fill=GOLD); s.T(rx+28,422+i*42,t,"Regular",15,INK)
    skyline(s,MX,655,(210,225,213))   # 공공건축물 스카이라인 (하단 모티프)
    footer(s,13,"출처: 공공기관 소유 어린이집 BIM · 현장 데이터")
    s.save("slide-13.png")

def s14():
    s=Slide(WHITE)
    title_block(s,"04 · 실증 사례","전체 보강 시 대상 건물은 인증 미달에서 ZEB 4등급으로 도약한다")
    # before / after
    by=300
    s.rect(MX,by,MX+360,by+200,fill=TINT,r=12)
    s.T(MX+30,by+30,"현재","Medium",15,GRAY)
    s.T(MX+30,by+58,"인증 미달","Bold",34,INK)
    s.T(MX+30,by+118,"1차에너지소요량","Regular",14,GRAY)
    s.T(MX+30,by+142,"166.7","Bold",30,GRAY)
    s.T(MX+150,by+150,"kWh/㎡·년","Regular",12,GRAY_L)
    arrow(s,MX+380,MX+470,by+100,GOLD,2.4)
    ax=MX+490
    s.rect(ax,by,ax+360,by+200,fill=(232,242,234),r=12)
    s.T(ax+30,by+30,"권장 전체 보강 후","Medium",15,GREEN)
    s.T(ax+30,by+58,"ZEB 4등급","Bold",34,GREEN)
    s.T(ax+30,by+118,"1차에너지소요량","Regular",14,GRAY)
    s.T(ax+30,by+142,"66.0","Bold",30,GREEN)
    s.T(ax+130,by+150,"< 90 기준 ✓","Medium",13,GOLD)
    # 등급 근거를 시각화 — 자립률 곡선이 아니라 소요량이 기준선 아래로 내려가는 그림
    mx0=MX+900; mx1=W-MX; my0=by+40; my1=by+150
    s.T(mx0,by+18,"제2호 · 소요량 하강","Regular",11.5,GRAY_L)
    for i in range(26):
        t=i/25; x=mx0+t*(mx1-mx0); yv=my0+(t**1.6)*(my1-my0)
        col=GRAY_L if t<0.5 else GREEN
        s.circle(x,yv,2.4,fill=col)
    s.T(mx0,by+168,"4등급 기준선 90 통과","Regular",11.5,GREEN)
    s.T(MX,540,"등급은 제1호(자립률) 또는 제2호(소요량) 중 상위로 정한다. 대상 건물은 태양광이 없어 자립률 0% →","Regular",13.5,GRAY)
    s.T(MX,562,"4등급은 태양광이 아니라 외피·설비 효율개선으로 소요량을 90 아래로 낮춰서 나온다. (BEMS 설치 가정)","Regular",13.5,GRAY)
    footer(s,14,"산정: ZEB-ROI 엔진 · 제2호 근거 (절감률 67%는 가정치 · 임계 55%)")
    s.save("slide-14.png")

def s15():
    s=Slide(WHITE)
    title_block(s,"04 · 실증 사례","자부담 0.91억 투자는 NPV +1.08억·IRR 14.7%로 회수된다")
    s.T(MX,150,"외피 보강 기준 경제성 · 할인회수 8.2년 · ZEB 4등급은 전체 보강(5.31억) 시 달성","Regular",13,GRAY)
    # Panel A — 사업비 흐름 (그룹 박스)
    pax,paw=MX,600
    s.rect(pax,196,pax+paw,330,fill=TINT,r=10)
    s.T(pax+24,216,"외피 보강 사업비 흐름","Bold",15,INK)
    flow=[("1.81억","Max Cost",GREEN),("-0.90억","정부 보조 50%",GOLD),("0.91억","자부담",INK)]
    for i,(val,lab,col) in enumerate(flow):
        fx=pax+24+i*208
        s.T(fx,256,val,"Bold",27,col); s.T(fx,294,lab,"Regular",12.5,GRAY)
        if i<2: arrow(s,fx+138,fx+186,272,GRAY_L,1.5)
    # Panel B — 절감·회수 (그룹 박스)
    pbx=pax+paw+24
    s.rect(pbx,196,W-MX,330,fill=TINT,r=10)
    s.T(pbx+24,216,"에너지 절감 · 회수","Bold",15,INK)
    s.T(pbx+24,256,"1,238만원","Bold",27,GREEN); s.T(pbx+24,294,"연간 절감","Regular",12.5,GRAY)
    s.T(pbx+250,256,"7.3년","Bold",27,INK); s.T(pbx+250,294,"단순 회수","Regular",12.5,GRAY)
    s.line(MX,360,W-MX,360,LINE,0.8)
    s.T(MX,384,"현금흐름 경제성  ·  자부담 기준 20년 · 할인율 4.5%","Bold",16,INK)
    cards=[("NPV","+1.08억","순현재가치",GREEN),("IRR","14.7%","내부수익률",GREEN),
           ("B–C","2.19","편익/비용",GREEN),("할인회수","8.2년","자부담 대비",INK)]
    cw=(W-2*MX-3*24)/4
    for i,(t,v,sub,col) in enumerate(cards):
        x=MX+i*(cw+24); s.rect(x,426,x+cw,596,fill=TINT,r=10)
        s.T(x+24,452,t,"Medium",14,GRAY); s.T(x+24,486,v,"Bold",38,col); s.T(x+24,554,sub,"Regular",12.5,GRAY)
    s.T(MX,616,"수익환원 자산가치 2.48억(환원율 5%)은 같은 절감의 다른 관점으로, NPV와 합산하지 않는다.","Regular",12.5,GRAY)
    footer(s,15,"출처: 조달청 07 단가DB·08 간접공사비 / 산정: ZEB-ROI 엔진")
    s.save("slide-15.png")

def s16():
    s=Slide(GREEN_DK)
    s.T(MX,62,"05 · 전망","Bold",14,GOLD,track=2.0)
    s.T(MX,92,"BIM→ZEB→경제성을 잇는 의사결정 자동화,","Bold",TITLE_SZ,WHITE)
    s.T(MX,92+s.lh("Bold",TITLE_SZ)*1.12,"그리고 그 다음은 AI다","Bold",TITLE_SZ,WHITE)
    # left 기대효과
    s.T(MX,236,"기대효과","Bold",16,GOLD)
    for i,t in enumerate(["수일~수주 걸리던 계산을 즉시 통합","단가DB·법령 근거 기반의 신뢰성","공공건축물 그린리모델링 확산 가속"]):
        s.circle(MX+8,276+i*48+9,4,fill=(150,200,160)); s.T(MX+28,276+i*48,t,"Regular",16,(210,222,213))
    # right AI 방향
    rx=MX+560
    s.T(rx,236,"앞으로의 AI 발전방향","Bold",16,GOLD)
    for i,t in enumerate(["BIM ↔ 에너지 시뮬레이션 실시간 연계","멀티에이전트가 설계대안 자동 생성·평가","인증·보조금 심사 자동화"]):
        token(s,rx+12,276+i*48+9,str(i+1),(40,80,48),13); s.T(rx+36,276+i*48,t,"Regular",16,(210,222,213))
    s.line(MX,470,W-MX,470,(58,78,64),0.8)
    s.T(MX,492,"남은 과제","Medium",14,(150,165,156))
    s.T(MX,516,"근거기반(RAG)으로 환각 방지 · 데이터 표준화(IFC) · AI 산정 결과의 책임소재","Regular",14.5,(190,205,193))
    s.T(MX,600,"BIM 한 번으로 ZEB 의사결정을 일원화한다 — 그린리모델링 확산의 디지털 기반.","Bold",18,WHITE)
    s.T(W-MX,648,"감사합니다  ·  Q & A","Medium",15,GOLD,anchor="ra")
    s.save("slide-16.png")

def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    blank=prs.slide_layouts[6]
    for i in range(1,17):
        sl=prs.slides.add_slide(blank)
        sl.shapes.add_picture(SLIDES+f"slide-{i:02d}.png",0,0,width=prs.slide_width,height=prs.slide_height)
    out=BASE+"ZEB-ROI_졸업설계_발표.pptx"; prs.save(out); print("PPTX:",out)

if __name__=="__main__":
    for fn in [s01,s02,s03,s04,s05,s06,s07,s08,s09,s10,s11,s12,s13,s14,s15,s16]:
        fn()
    print("16 slides rendered")
    build_pptx()
