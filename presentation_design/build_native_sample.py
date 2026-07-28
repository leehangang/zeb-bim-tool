# -*- coding: utf-8 -*-
"""
ZEB-ROI 발표 — 네이티브(편집 가능) PPTX 샘플 3장.
python-pptx로 텍스트박스·도형을 직접 생성 → PowerPoint에서 글자 수정 가능.
좌표는 HTML 덱(1280x720px)을 그대로 따름: 1px = 9525 EMU.
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

BASE = "C:/Users/이혁주/Desktop/zeb-chatbot/presentation_design/"
IMG = BASE + "images/"
EMU = 9525
def P(px): return Emu(int(round(px*EMU)))
def PT(px): return Pt(px*0.75)   # 96dpi px -> pt

GREEN=RGBColor(0x1B,0x5E,0x20); GREEN_DK=RGBColor(0x0C,0x28,0x15)
GOLD=RGBColor(0xC1,0x8A,0x2D); INK=RGBColor(0x1A,0x20,0x1C)
GRAY=RGBColor(0x5C,0x66,0x5F); GRAY_L=RGBColor(0x9B,0xA3,0x9C)
LINE=RGBColor(0xE3,0xE6,0xE2); TINT=RGBColor(0xEE,0xF3,0xEF)
GOLD_SOFT=RGBColor(0xF4,0xEF,0xE2); WHITE=RGBColor(0xFF,0xFF,0xFF)
SOFT=RGBColor(0xE6,0xF0,0xE8); T_EYE=RGBColor(0x96,0xC8,0xA0); SUBC=RGBColor(0xBE,0xCD,0xC1)
FONT="Noto Sans KR"

prs = Presentation()
prs.slide_width = P(1280); prs.slide_height = P(720)
BLANK = prs.slide_layouts[6]

_CONV = BASE + "images/_pptx/"
os.makedirs(_CONV, exist_ok=True)
def safe_img(path):
    """python-pptx는 WEBP 미지원 → PNG로 변환한 사본 경로 반환."""
    try:
        im = Image.open(path)
        if im.format in ("JPEG","PNG","GIF","BMP","TIFF"):
            return path
        out = _CONV + os.path.splitext(os.path.basename(path))[0] + ".png"
        im.convert("RGB").save(out, "PNG")
        return out
    except Exception:
        return path

def _setfont(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea","a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)

def tb(slide,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,spacing=1.0):
    """runs: list of paragraphs; each paragraph = list of (text,size,bold,color)."""
    box=slide.shapes.add_textbox(P(x),P(y),P(w),P(h)); tf=box.text_frame
    tf.word_wrap=True
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    tf.vertical_anchor=anchor
    for i,para in enumerate(runs):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align
        if spacing: p.line_spacing=spacing
        for (txt,size,bold,color) in para:
            r=p.add_run(); r.text=txt; r.font.size=PT(size); r.font.bold=bold
            r.font.color.rgb=color; _setfont(r)
    return box

def rect(slide,x,y,w,h,fill=None,line=None,lw=1.0,round=False,radius=0.10):
    shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,
                               P(x),P(y),P(w),P(h))
    if round:
        try: shp.adjustments[0]=radius
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=PT(lw)
    shp.shadow.inherit=False
    return shp

def oval(slide,x,y,d,fill,num=None,numcolor=WHITE,numsize=11):
    shp=slide.shapes.add_shape(MSO_SHAPE.OVAL,P(x),P(y),P(d),P(d))
    shp.fill.solid(); shp.fill.fore_color.rgb=fill; shp.line.fill.background()
    shp.shadow.inherit=False
    if num is not None:
        tf=shp.text_frame; tf.word_wrap=False
        tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
        tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=num; r.font.size=PT(numsize); r.font.bold=True
        r.font.color.rgb=numcolor; _setfont(r)
    return shp

def hline(slide,x,y,w,color,wt=0.8):
    cn=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,P(x),P(y),P(x+w),P(y))
    cn.line.color.rgb=color; cn.line.width=PT(wt); cn.shadow.inherit=False
    return cn

def footer(slide,page,source,dark=False):
    c = RGBColor(0x8C,0x9C,0x92) if dark else GRAY_L
    hline(slide,88,678,1104, RGBColor(0x3A,0x4E,0x40) if dark else LINE, 0.8)
    tb(slide,88,684,520,20,[[(f"{page:02d}",11,True,c),("   ZEB-ROI · 그린리모델링 의사결정 플랫폼",11,False,c)]])
    tb(slide,672,684,520,20,[[(source,10.5,False,c)]],align=PP_ALIGN.RIGHT)

# ===================================================================
# SLIDE 1 — TITLE (dark)
# ===================================================================
def s_title():
    s=prs.slides.add_slide(BLANK)
    rect(s,0,0,1280,720,fill=GREEN_DK)
    tb(s,88,150,900,22,[[("ZERO ENERGY BUILDING · GREEN REMODELING",10.5,False,T_EYE)]])
    tb(s,88,190,1100,140,[
        [("BIM 한 번으로 ZEB 등급을 평가하고",34.5,True,WHITE)],
        [("그린리모델링 전 과정을 설계하다",34.5,True,WHITE)],
    ],spacing=1.18)
    tb(s,88,338,1000,30,[[("제로에너지건축물(ZEB) 의사결정 통합 플랫폼",14,False,SUBC)]])
    hline(s,88,604,360,RGBColor(0x46,0x60,0x4C),1.0)
    tb(s,88,616,500,24,[[("졸업설계 작품 · 2026",10.5,False,RGBColor(0xAF,0xC2,0xB2))]])
    tb(s,692,616,500,24,[[("사례 · 공공기관 소유 어린이집",10,False,GOLD)]],align=PP_ALIGN.RIGHT)

# ===================================================================
# SLIDE 2 — MODE 01 (screenshot left + editable explanation right)
# ===================================================================
def s_mode():
    s=prs.slides.add_slide(BLANK)
    tb(s,88,60,900,20,[[("03 · 제안 · MODE 01",10.5,True,GOLD)]])
    oval(s,88,84,40,GREEN,num="01",numsize=12)
    tb(s,140,84,1050,44,[[("🏢 BIM 진단 + ZEB 등급 평가 — 플랫폼의 핵심 엔진",22,True,INK)]],
       anchor=MSO_ANCHOR.MIDDLE)
    # screenshot left (contain: 620 wide -> 421 tall)
    rect(s,88,196,624,428,fill=RGBColor(0xF4,0xF6,0xF4),line=LINE,lw=1.0,round=True,radius=0.03)
    p=IMG+"app-bim.jpg"
    if os.path.exists(p): s.shapes.add_picture(safe_img(p),P(90),P(199),P(620),P(421))
    # right explanation
    rx=740
    tb(s,rx,200,452,52,[[("BIM 한 번 업로드로 ZEB 인증 등급과 보강 계획을 동시에 산출한다.",11.5,True,INK)]],spacing=1.15)
    steps=[("1","Dynamo/Revit 추출 BIM JSON 업로드 (또는 데모 3종)"),
           ("2","객체를 11개 GR 기술요소로 자동 매핑"),
           ("3","두 잣대로 동시 채점 — 정량평가 100점(A+~D) + ZEB 인증 등급"),
           ("4","06 별표1 지역별 열관류율 적합성 자동 판정")]
    yy=266
    for n,txt in steps:
        oval(s,rx,yy,26,SOFT,num=n,numcolor=GREEN,numsize=10)
        tb(s,rx+38,yy-3,418,44,[[(txt,10.5,False,INK)]],spacing=1.12,anchor=MSO_ANCHOR.TOP)
        yy+=60
    rect(s,rx,514,452,62,fill=TINT,round=True,radius=0.12)
    tb(s,rx+15,524,422,44,[[("핵심 기능",9.5,True,GREEN),
        ("  ·  진단 / ROI 보강계획 / 최적화 3개 탭 · 100점 항목별 분해 · PDF 진단 리포트",9.5,False,GRAY)]],spacing=1.2)
    # io tags
    rect(s,rx,590,196,30,fill=WHITE,line=LINE,lw=1.0,round=True,radius=0.5)
    tb(s,rx,597,196,18,[[("입력 · Revit/Dynamo JSON",9,False,GRAY)]],align=PP_ALIGN.CENTER)
    rect(s,rx+208,590,196,30,fill=WHITE,line=LINE,lw=1.0,round=True,radius=0.5)
    tb(s,rx+208,597,196,18,[[("출력 · 점수·ZEB등급·PDF",9,False,GRAY)]],align=PP_ALIGN.CENTER)
    footer(s,12,"엔진: core.bim_diagnoser + core.zeb_evaluator · 01 정량평가표 · 06 별표1")

# ===================================================================
# SLIDE 3 — 핵심 방법론 (infographic: 요건 박스 + 등급 사다리)
# ===================================================================
def s_method():
    s=prs.slides.add_slide(BLANK)
    tb(s,88,60,900,20,[[("03 · 핵심 방법론",10.5,True,GOLD)]])
    tb(s,88,84,1110,80,[[("ZEB 인증은 자립률 또는 1차에너지소요량, 그리고 BEMS로 결정된다",22,True,INK)]],spacing=1.12)
    # requirement row
    y=246
    def reqbox(x,h_,d_):
        rect(s,x,y,238,86,fill=TINT,round=True,radius=0.12)
        tb(s,x+22,y+18,200,20,[[(h_,10.5,True,GREEN)]])
        tb(s,x+22,y+44,200,24,[[(d_,12,False,INK)]])
    reqbox(88,"제1호","에너지 자립률")
    tb(s,330,y+30,52,26,[[("OR",11.5,True,GOLD)]],align=PP_ALIGN.CENTER)
    reqbox(380,"제2호","1차에너지소요량")
    tb(s,624,y+30,56,26,[[("AND",11.5,True,GOLD)]],align=PP_ALIGN.CENTER)
    reqbox(684,"제3호","BEMS 설치")
    tb(s,936,y+28,40,30,[[("→",16,False,GRAY_L)]],align=PP_ALIGN.CENTER)
    tb(s,986,y+30,206,26,[[("인증 등급",13.5,True,INK)]])
    # ladder
    tb(s,88,398,700,24,[[("ZEB 등급 — 에너지 자립률 기준(제1호)",12.5,True,INK)]])
    grades=[("+등급","120%",GREEN),("1등급","100%",GREEN),("2등급","80%",GREEN),
            ("3등급","60%",GOLD),("4등급","40%",GOLD),("5등급","20%",GOLD)]
    gw=(1104-5*16)/6
    gy=434
    for i,(g,v,col) in enumerate(grades):
        gx=88+i*(gw+16)
        rect(s,gx,gy,gw,92,fill=TINT,round=True,radius=0.10)
        tb(s,gx,gy+22,gw,24,[[(g,14,True,col)]],align=PP_ALIGN.CENTER)
        tb(s,gx,gy+54,gw,20,[[(v,11.5,False,GRAY)]],align=PP_ALIGN.CENTER)
    tb(s,88,560,1110,20,[[("전력 1차에너지 환산계수 × 2.75 적용  ·  인증등급은 제1·2호 중 상위 등급으로 산정",10,False,GRAY)]])
    footer(s,17,"출처: 03 ZEB 인증기준 고시 · 건축물 에너지효율등급 기준")

s_title(); s_mode(); s_method()
out=BASE+"ZEB-ROI_네이티브샘플.pptx"
prs.save(out)
print("saved:",out,"| slides:",len(prs.slides._sldIdLst))
